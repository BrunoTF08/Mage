from pathlib import Path
import shutil

from PIL import Image, ImageEnhance, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\BRUNO\Documents\ProjetoTwin\Mage")
OUT_DIR = ROOT / "Presentation"
PPTX_PATH = OUT_DIR / "Mage_Apresentacao_Projeto.pptx"
BACKUP_PATH = OUT_DIR / "Mage_Apresentacao_Projeto_original.pptx"

IMG_SCENE = OUT_DIR / "Mage_SceneView.png"
IMG_GAME = OUT_DIR / "Mage_GameView.png"
IMG_PLAY = OUT_DIR / "Mage_PlayMode.png"
IMG_HEALTH = ROOT / "Assets" / "Wizard" / "barra de vida original 576x256.png"
IMG_CROSSHAIR_SMALL = (
    ROOT
    / "Assets"
    / "Wizard"
    / "Scripts"
    / "mira"
    / "Crosshair Pack uno.png"
)


NAVY = RGBColor(8, 13, 34)
NAVY2 = RGBColor(13, 21, 50)
PANEL = RGBColor(19, 31, 67)
PANEL2 = RGBColor(27, 41, 83)
CYAN = RGBColor(77, 210, 255)
BLUE = RGBColor(61, 122, 255)
GOLD = RGBColor(255, 188, 73)
GREEN = RGBColor(83, 214, 132)
RED = RGBColor(255, 104, 104)
WHITE = RGBColor(246, 248, 255)
MUTED = RGBColor(184, 198, 232)
GRAY = RGBColor(122, 139, 180)
INK = RGBColor(5, 9, 24)
FONT = "Aptos"
FONT_HEAD = "Aptos Display"
FONT_CODE = "Consolas"


def make_overlay(src: Path, dst: Path, darken: float = 0.58, blur: int = 0) -> None:
    img = Image.open(src).convert("RGB")
    if blur:
        img = img.filter(ImageFilter.GaussianBlur(blur))
    img = ImageEnhance.Brightness(img).enhance(darken)
    img.save(dst, quality=95)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


prs = new_prs()
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height
slide_no = 0


def add_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_image(slide, path: Path, x, y, w, h):
    return slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
    )


def add_footer(slide, label="Mage | Projeto Unity"):
    global slide_no
    slide_no += 1
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(4.2), Inches(0.22))
    tf = tx.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = FONT
    p.font.size = Pt(8.5)
    p.font.color.rgb = GRAY

    tx2 = slide.shapes.add_textbox(Inches(12.35), Inches(7.08), Inches(0.45), Inches(0.22))
    tf2 = tx2.text_frame
    tf2.margin_left = 0
    tf2.margin_right = 0
    tf2.margin_top = 0
    tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.text = str(slide_no)
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.name = FONT
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = GRAY


def add_title(slide, title, subtitle=None, y=0.45, color=WHITE):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(y), Inches(9.2), Inches(0.75))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_HEAD
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.68), Inches(y + 0.7), Inches(9.0), Inches(0.35)
        )
        tf2 = sub.text_frame
        tf2.margin_left = 0
        tf2.margin_right = 0
        tf2.margin_top = 0
        tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.name = FONT
        p2.font.size = Pt(13)
        p2.font.color.rgb = MUTED


def add_tag(slide, text, x, y, w=1.9, color=CYAN):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.adjustments[0] = 0.12
    tf = shp.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = NAVY
    return shp


def add_card(slide, x, y, w, h, title=None, body=None, accent=CYAN, fill=PANEL):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = RGBColor(42, 59, 107)
    card.line.width = Pt(1)
    card.adjustments[0] = 0.06

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    if title:
        tb = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.18), Inches(w - 0.45), Inches(0.34))
        tf = tb.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE

    if body:
        bb = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.58), Inches(w - 0.45), Inches(h - 0.68))
        tf = bb.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        if isinstance(body, list):
            for i, item in enumerate(body):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item
                p.font.name = FONT
                p.font.size = Pt(10.7)
                p.font.color.rgb = MUTED
                p.space_after = Pt(4)
        else:
            p = tf.paragraphs[0]
            p.text = body
            p.font.name = FONT
            p.font.size = Pt(11)
            p.font.color.rgb = MUTED
            p.line_spacing = 1.08
    return card


def add_metric(slide, value, label, x, y, w, accent=CYAN):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.78))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = value
    p.font.name = FONT_HEAD
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = accent

    lb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.62), Inches(w), Inches(0.35))
    tf2 = lb.text_frame
    tf2.margin_left = 0
    tf2.margin_right = 0
    tf2.margin_top = 0
    tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.name = FONT
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = MUTED


def add_code(slide, code, x, y, w, h, title=None):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(2, 6, 18)
    box.line.color.rgb = RGBColor(58, 80, 140)
    box.adjustments[0] = 0.04

    y_offset = 0.18
    if title:
        tb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.17), Inches(w - 0.4), Inches(0.3))
        tf = tb.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = GOLD
        y_offset = 0.55

    tb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + y_offset), Inches(w - 0.4), Inches(h - y_offset - 0.15))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = FONT_CODE
    p.font.size = Pt(8.5)
    p.font.color.rgb = RGBColor(215, 230, 255)
    return box


def add_section_label(slide, text, x, y, color=GOLD):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.0), Inches(0.25))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.name = FONT
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = color
    return tb


def prepare_images():
    make_overlay(IMG_GAME, OUT_DIR / "Mage_GameView_Dark.jpg", 0.55, 0)
    make_overlay(IMG_SCENE, OUT_DIR / "Mage_SceneView_Dark.jpg", 0.62, 0)
    make_overlay(IMG_PLAY, OUT_DIR / "Mage_PlayMode_Dark.jpg", 0.52, 0)
    make_overlay(IMG_PLAY, OUT_DIR / "Mage_PlayMode_Blur.jpg", 0.40, 6)


def build_slides():
    # 1
    slide = prs.slides.add_slide(BLANK)
    add_image(slide, OUT_DIR / "Mage_PlayMode_Dark.jpg", 0, 0, 13.333, 7.5)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = INK
    overlay.fill.transparency = 35
    overlay.line.fill.background()
    add_tag(slide, "UNITY 2022.3 | URP", 0.75, 0.65, 2.25, CYAN)
    title = slide.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(7.1), Inches(1.5))
    tf = title.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "Mage"
    p.font.name = FONT_HEAD
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.85), Inches(7.3), Inches(0.75))
    tf = sub.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "Prot\u00f3tipo de combate m\u00e1gico em terceira pessoa"
    p.font.name = FONT
    p.font.size = Pt(20)
    p.font.color.rgb = MUTED
    add_card(
        slide,
        0.78,
        4.78,
        4.5,
        1.22,
        "Base da apresenta\u00e7\u00e3o",
        [
            "Captura real em Play Mode",
            "Explica\u00e7\u00e3o dos scripts principais",
            "Resumo das anima\u00e7\u00f5es do mago",
        ],
        CYAN,
        RGBColor(9, 18, 45),
    )
    add_footer(slide, "Mage | Apresenta\u00e7\u00e3o do projeto")

    # 2
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_title(slide, "Vis\u00e3o Geral", "O que o projeto j\u00e1 comunica hoje")
    add_card(
        slide,
        0.7,
        1.65,
        3.6,
        3.75,
        "Conceito",
        "Mage \u00e9 um prot\u00f3tipo 3D de fantasia focado em um mago jog\u00e1vel, c\u00e2mera de terceira pessoa e combate por proj\u00e9teis m\u00e1gicos.",
        CYAN,
    )
    add_card(
        slide,
        4.65,
        1.65,
        3.6,
        3.75,
        "Experi\u00eancia",
        "O jogador controla o mago em uma arena de teste, mira com c\u00e2mera sobre o ombro, ataca com magia e recebe feedback por anima\u00e7\u00e3o, HUD e \u00e1udio.",
        GOLD,
    )
    add_card(
        slide,
        8.6,
        1.65,
        3.6,
        3.75,
        "Estado atual",
        "A base de gameplay est\u00e1 montada: movimento, corrida, mira, vida, barra de vida, disparo m\u00e1gico, som e cena jog\u00e1vel de valida\u00e7\u00e3o.",
        GREEN,
    )
    add_metric(slide, "3.0", "vers\u00e3o do produto", 0.85, 5.95, 2.6, CYAN)
    add_metric(slide, "3039", "assets no projeto", 4.0, 5.95, 2.6, GOLD)
    add_metric(slide, "URP", "pipeline gr\u00e1fico ativo", 7.15, 5.95, 2.6, GREEN)
    add_metric(slide, "Windows", "target atual", 10.3, 5.95, 2.6, BLUE)
    add_footer(slide)

    # 3
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, NAVY2)
    add_title(slide, "Imagem em Play Mode", "Captura do jogo rodando no Editor")
    add_image(slide, IMG_PLAY, 0.72, 1.35, 7.55, 4.25)
    add_card(
        slide,
        8.55,
        1.35,
        3.85,
        1.15,
        "O que a imagem mostra",
        "Mago enquadrado pela c\u00e2mera em terceira pessoa, arena de teste e dire\u00e7\u00e3o de mira \u00e0 frente.",
        CYAN,
    )
    add_card(
        slide,
        8.55,
        2.8,
        3.85,
        1.15,
        "Por que \u00e9 importante",
        "A captura confirma que a cena n\u00e3o \u00e9 apenas montagem no Editor: ela executa em Play Mode.",
        GOLD,
    )
    add_card(
        slide,
        8.55,
        4.25,
        3.85,
        1.15,
        "Valida\u00e7\u00e3o",
        "Imagem renderizada pela Main Camera durante Play Mode via Unity MCP.",
        GREEN,
    )
    add_footer(slide)

    # 4
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, NAVY2)
    add_title(slide, "Loop de Gameplay", "A base jog\u00e1vel que aparece nos scripts e na cena")
    steps = [
        ("1", "Explorar", "Movimento WASD relativo \u00e0 c\u00e2mera e corrida com Shift."),
        ("2", "Mirar", "C\u00e2mera sobre o ombro, cursor travado e alvo atualizado \u00e0 frente."),
        ("3", "Conjurar", "Clique esquerdo aciona anima\u00e7\u00e3o, som e proj\u00e9til m\u00e1gico."),
        ("4", "Reagir", "Vida muda no PlayerHealth e o HUD troca os sprites da barra."),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.75 + i * 3.1
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(1.72), Inches(0.62), Inches(0.62))
        circ.fill.solid()
        circ.fill.fore_color.rgb = [CYAN, GOLD, BLUE, GREEN][i]
        circ.line.fill.background()
        tf = circ.text_frame
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = num
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        add_card(slide, x, 2.55, 2.6, 2.15, title, body, [CYAN, GOLD, BLUE, GREEN][i], PANEL)
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.35), Inches(1.9), Inches(0.65), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(51, 72, 125)
            arrow.line.fill.background()
    add_card(
        slide,
        1.35,
        5.45,
        10.65,
        0.85,
        "Leitura de design",
        "A for\u00e7a do projeto est\u00e1 em um prot\u00f3tipo interativo: c\u00e2mera, personagem, mira, anima\u00e7\u00e3o e feedback est\u00e3o conectados.",
        CYAN,
        RGBColor(12, 26, 57),
    )
    add_footer(slide)

    # 5
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_title(slide, "Programa\u00e7\u00e3o: Scripts Principais", "Responsabilidade de cada componente do jogador")
    add_card(slide, 0.75, 1.35, 2.85, 2.0, "PlayerController", ["Movimento e corrida", "Ataque por clique", "Raycast de mira", "Instancia proj\u00e9til"], CYAN)
    add_card(slide, 3.9, 1.35, 2.85, 2.0, "ShoulderCameraController", ["Yaw/pitch do mouse", "Offset no ombro", "Suaviza\u00e7\u00e3o", "Cursor travado"], GOLD)
    add_card(slide, 7.05, 1.35, 2.85, 2.0, "MagicProjectile", ["Velocidade", "Tempo de vida", "Seguimento do alvo", "Colis\u00e3o"], GREEN)
    add_card(slide, 10.2, 1.35, 2.85, 2.0, "PlayerHealth", ["Vida m\u00e1xima", "Dano e cura", "Evento HealthChanged", "Percentual de vida"], BLUE)
    add_card(slide, 1.55, 4.0, 3.1, 1.55, "HealthBarUI", "Escuta o evento de vida e troca o sprite da barra conforme o percentual atual.", GREEN)
    add_card(slide, 5.15, 4.0, 3.1, 1.55, "PlayerFootsteps", "Sincroniza passos com estados de caminhada/corrida e varia volume/pitch.", GOLD)
    add_card(slide, 8.75, 4.0, 3.1, 1.55, "AimTarget", "Mant\u00e9m um ponto \u00e0 frente da c\u00e2mera para orientar mira e proj\u00e9teis.", CYAN)
    add_footer(slide)

    # 6
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, NAVY2)
    add_title(slide, "Programa\u00e7\u00e3o: Ataque M\u00e1gico", "Como o clique vira um proj\u00e9til")
    add_code(
        slide,
        "Input.GetMouseButtonDown(0)\n"
        "SetAnimatorTrigger(AttackHash)\n"
        "QueueMagicAttack()\n"
        "CastMagicAttack()\n"
        "SpawnMagicProjectile()\n"
        "projectile.Launch(aimDirection, aimTarget)",
        0.75,
        1.45,
        5.5,
        3.05,
        "Fluxo no PlayerController.cs",
    )
    add_card(
        slide,
        6.65,
        1.45,
        2.75,
        1.28,
        "1. Entrada",
        "O clique esquerdo inicia o ataque se o mago n\u00e3o estiver atacando.",
        CYAN,
    )
    add_card(
        slide,
        9.75,
        1.45,
        2.75,
        1.28,
        "2. Anima\u00e7\u00e3o",
        "Triggers Attack e IsAttack chamam o estado Attack04.",
        GOLD,
    )
    add_card(
        slide,
        6.65,
        3.08,
        2.75,
        1.28,
        "3. Mira",
        "Um raycast da c\u00e2mera define o ponto de impacto esperado.",
        BLUE,
    )
    add_card(
        slide,
        9.75,
        3.08,
        2.75,
        1.28,
        "4. Proj\u00e9til",
        "O prefab \u00e9 criado no SpellFirePoint e segue o AimTarget.",
        GREEN,
    )
    add_card(
        slide,
        1.1,
        5.25,
        11.2,
        0.78,
        "Par\u00e2metros atuais",
        "Velocidade do proj\u00e9til: 25 | tempo de vida: 3s | atraso de conjura\u00e7\u00e3o: 0,18s | raycast de mira: 80 unidades.",
        CYAN,
        RGBColor(12, 26, 57),
    )
    add_footer(slide)

    # 7
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_title(slide, "Programa\u00e7\u00e3o: C\u00e2mera e Mira", "Controle de terceira pessoa")
    add_image(slide, OUT_DIR / "Mage_SceneView_Dark.jpg", 0.65, 1.35, 5.3, 3.0)
    add_code(
        slide,
        "yaw += Input.GetAxisRaw(\"Mouse X\") * mouseSensitivity;\n"
        "pitch -= Input.GetAxisRaw(\"Mouse Y\") * mouseSensitivity;\n"
        "pivotPosition = player.position + yawRotation * shoulderOffset;\n"
        "cameraPosition = pivotPosition - cameraRotation * Vector3.forward * distance;\n"
        "player.rotation = Quaternion.Slerp(...);",
        6.25,
        1.35,
        6.2,
        2.45,
        "ShoulderCameraController.cs",
    )
    add_card(
        slide,
        6.25,
        4.2,
        2.95,
        1.32,
        "Valores na cena",
        "Dist\u00e2ncia 3,2 | sensibilidade 2,2 | pitch entre -30 e 65 graus.",
        CYAN,
    )
    add_card(
        slide,
        9.55,
        4.2,
        2.95,
        1.32,
        "AimTarget",
        "O alvo fica \u00e0 frente da Main Camera e guia a dire\u00e7\u00e3o da magia.",
        GOLD,
    )
    add_footer(slide)

    # 8
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, NAVY2)
    add_title(slide, "Programa\u00e7\u00e3o: Vida, HUD e \u00c1udio", "Feedback para o jogador")
    add_code(
        slide,
        "public event Action<float, float> HealthChanged;\n"
        "SetHealth(currentHealth - amount);\n"
        "HealthChanged?.Invoke(currentHealth, maxHealth);\n\n"
        "percent = currentHealth / maxHealth;\n"
        "image.sprite = healthSprites[spriteIndex];",
        0.75,
        1.45,
        5.45,
        2.9,
        "PlayerHealth.cs + HealthBarUI.cs",
    )
    add_image(slide, IMG_HEALTH, 6.65, 1.75, 5.55, 2.47)
    add_card(
        slide,
        0.95,
        5.0,
        3.55,
        1.0,
        "Vida",
        "Dano, cura e clamp impedem valores fora de 0 a 100.",
        GREEN,
    )
    add_card(
        slide,
        4.9,
        5.0,
        3.55,
        1.0,
        "HUD",
        "A barra visual muda de sprite de acordo com a vida.",
        CYAN,
    )
    add_card(
        slide,
        8.85,
        5.0,
        3.55,
        1.0,
        "\u00c1udio",
        "Passos e magia usam AudioSource e AudioClip.",
        GOLD,
    )
    add_footer(slide)

    # 9
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_title(slide, "Anima\u00e7\u00f5es do Mago", "Estados e par\u00e2metros conectados ao gameplay")
    add_card(
        slide,
        0.75,
        1.35,
        3.65,
        4.1,
        "Controller ativo",
        [
            "WizardAnimControl",
            "Avatar: WizardBodyMeshAvatar",
            "Controller em Assets/WizardPBR/Animations",
            "Estados usados pelo PlayerController",
        ],
        CYAN,
    )
    add_card(
        slide,
        4.85,
        1.35,
        3.65,
        4.1,
        "Par\u00e2metros",
        [
            "Attack: Trigger",
            "IsAttack: Trigger",
            "Walk, WalkBack: Bool",
            "IsWalkingRight/Left: Bool",
            "IsRunning: Bool",
            "IsJump: Bool",
        ],
        GOLD,
    )
    add_card(
        slide,
        8.95,
        1.35,
        3.65,
        4.1,
        "Estados principais",
        [
            "Idle03",
            "BattleWalkForward",
            "BattleWalkBack",
            "BattleWalkLeft/Right",
            "BattleRunForward",
            "Attack04",
        ],
        GREEN,
    )
    add_card(
        slide,
        1.3,
        5.9,
        10.6,
        0.72,
        "Como funciona",
        "O script n\u00e3o depende apenas de transi\u00e7\u00f5es do Animator: ele tamb\u00e9m usa CrossFadeInFixedTime para trocar de estado com blend curto.",
        BLUE,
        RGBColor(12, 26, 57),
    )
    add_footer(slide)

    # 10
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, NAVY2)
    add_title(slide, "Anima\u00e7\u00f5es e \u00c1udio de Passos", "Sincroniza\u00e7\u00e3o por estado de anima\u00e7\u00e3o")
    add_code(
        slide,
        "AnimatorStateInfo stateInfo = animator.GetCurrentAnimatorStateInfo(0);\n"
        "FootstepMode mode = GetFootstepMode(stateInfo);\n"
        "cycleTime = Mathf.Repeat(stateInfo.normalizedTime, 1f);\n"
        "if (cycleTime >= stepTimes.x) PlayFootstep(mode);\n"
        "if (cycleTime >= stepTimes.y) PlayFootstep(mode);",
        0.75,
        1.45,
        6.2,
        3.15,
        "PlayerFootsteps.cs",
    )
    add_card(slide, 7.35, 1.45, 4.95, 1.05, "Caminhada", "Estados BattleWalkForward, BattleWalkBack, BattleWalkRight e BattleWalkLeft.", CYAN)
    add_card(slide, 7.35, 2.78, 4.95, 1.05, "Corrida", "Estado BattleRunForward com volume maior e tempos de passo pr\u00f3prios.", GOLD)
    add_card(slide, 7.35, 4.1, 4.95, 1.05, "Varia\u00e7\u00e3o", "Pitch aleat\u00f3rio entre 0,92 e 1,08 evita som repetitivo.", GREEN)
    add_card(slide, 1.1, 5.45, 11.2, 0.78, "Resultado", "A movimenta\u00e7\u00e3o ganha resposta sonora conforme o estado atual do Animator e a velocidade horizontal do CharacterController.", BLUE, RGBColor(12, 26, 57))
    add_footer(slide)

    # 11
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_title(slide, "Cena Jog\u00e1vel", "Wizard Testes 1 como vitrine atual")
    add_image(slide, IMG_GAME, 0.7, 1.45, 7.25, 4.08)
    add_card(slide, 8.25, 1.45, 4.25, 1.15, "C\u00e2mera em terceira pessoa", "A Main Camera usa ShoulderCameraController e enquadra o mago por tr\u00e1s do ombro.", CYAN)
    add_card(slide, 8.25, 2.85, 4.25, 1.15, "Ambiente de teste", "Cen\u00e1rio noturno, Terrain, PP Volume, ilumina\u00e7\u00e3o e colisores de teste.", GOLD)
    add_card(slide, 8.25, 4.25, 4.25, 1.15, "HUD ativo", "A cena possui Canvas com CrossHair e HealthBar ligados ao fluxo de mira e vida.", GREEN)
    add_metric(slide, "94", "GameObjects na cena", 0.95, 5.95, 2.2, CYAN)
    add_metric(slide, "147", "componentes", 3.55, 5.95, 2.0, GOLD)
    add_metric(slide, "0", "refer\u00eancias quebradas", 5.9, 5.95, 2.4, GREEN)
    add_footer(slide)

    # 12
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, NAVY2)
    add_title(slide, "Personagem Jog\u00e1vel", "O mago como n\u00facleo do prot\u00f3tipo")
    add_image(slide, IMG_SCENE, 7.25, 1.15, 5.45, 3.07)
    add_card(slide, 0.75, 1.45, 5.95, 1.1, "Composi\u00e7\u00e3o em cena", "WizardStandardMaterial concentra Animator, PlayerController, Rigidbody, CapsuleCollider, CharacterController, PlayerHealth, AudioSource e PlayerFootsteps.", CYAN)
    add_card(slide, 0.75, 2.82, 5.95, 1.1, "Movimento", "Andar, recuar, mover lateralmente e correr usam CharacterController, gravidade e dire\u00e7\u00e3o relativa \u00e0 c\u00e2mera.", GOLD)
    add_card(slide, 0.75, 4.19, 5.95, 1.1, "Anima\u00e7\u00e3o", "Estados como Idle03, BattleWalkForward, BattleRunForward e Attack04 s\u00e3o acionados por par\u00e2metros e crossfades.", GREEN)
    add_card(slide, 7.25, 4.65, 5.45, 1.15, "Pontos de gameplay", "LookAt, SpellFirePoint e AimTarget d\u00e3o suporte ao disparo m\u00e1gico e \u00e0 mira din\u00e2mica.", BLUE)
    add_footer(slide)

    # 13
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_title(slide, "Arquitetura T\u00e9cnica", "Camadas principais do projeto")
    cols = [
        ("Gameplay", ["PlayerController", "PlayerHealth", "MagicProjectile", "PlayerFootsteps"], CYAN),
        ("C\u00e2mera", ["ShoulderCameraController", "AimTarget", "Cinemachine"], GOLD),
        ("UI", ["HealthBarUI", "Canvas", "CrossHair", "Sprites de vida"], GREEN),
        ("Editor", ["ShoulderCameraSceneSetup", "Unity MCP", "Build Settings"], BLUE),
    ]
    for i, (title, items, c) in enumerate(cols):
        add_card(slide, 0.75 + i * 3.08, 1.5, 2.65, 3.0, title, items, c, PANEL)
    add_card(slide, 0.95, 5.15, 2.65, 0.85, "Unity", "2022.3.22f1 LTS", CYAN, RGBColor(12, 26, 57))
    add_card(slide, 3.95, 5.15, 2.65, 0.85, "Render", "URP 14.0.10", GOLD, RGBColor(12, 26, 57))
    add_card(slide, 6.95, 5.15, 2.65, 0.85, "Input", "Input System 1.11.0 + eixos legados", GREEN, RGBColor(12, 26, 57))
    add_card(slide, 9.95, 5.15, 2.65, 0.85, "UI", "uGUI + TextMeshPro", BLUE, RGBColor(12, 26, 57))
    add_footer(slide)

    # 14
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, NAVY2)
    add_title(slide, "Status de Valida\u00e7\u00e3o", "Leitura feita no Editor e no projeto")
    status = [
        ("Compila\u00e7\u00e3o", "0 erros reportados pelo Unity MCP", GREEN),
        ("Console", "Sem erros recentes; apenas logs do MCP", GREEN),
        ("Cena ativa", "Wizard Testes 1 carregada", CYAN),
        ("Refer\u00eancias", "0 refer\u00eancias quebradas na cena ativa", GREEN),
    ]
    for i, (t, b, c) in enumerate(status):
        add_card(slide, 0.8 + (i % 2) * 6.05, 1.55 + (i // 2) * 1.6, 5.55, 1.15, t, b, c, PANEL)
    add_card(
        slide,
        0.8,
        5.05,
        11.6,
        1.05,
        "Ponto de aten\u00e7\u00e3o",
        "ProjectSettings/EditorBuildSettings.asset ainda lista Assets/Scenes/Wizard Testes.unity como cena habilitada, mas esse arquivo n\u00e3o foi encontrado em Assets/Scenes. Recomenda-se corrigir isso antes do build final.",
        GOLD,
        RGBColor(39, 34, 20),
    )
    add_footer(slide)

    # 15
    slide = prs.slides.add_slide(BLANK)
    add_image(slide, OUT_DIR / "Mage_PlayMode_Blur.jpg", 0, 0, 13.333, 7.5)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = INK
    overlay.fill.transparency = 18
    overlay.line.fill.background()
    add_title(slide, "Pr\u00f3ximos Passos", "Roteiro sugerido para transformar o prot\u00f3tipo em demo")
    roadmap = [
        ("1", "Ajustar build", "Remover ou corrigir a cena ausente nos Build Settings."),
        ("2", "Fechar combate", "Adicionar inimigos, dano recebido e dano causado pelo proj\u00e9til."),
        ("3", "Polir sensa\u00e7\u00e3o", "Melhorar VFX de magia, impacto, c\u00e2mera e mixagem de \u00e1udio."),
        ("4", "Estruturar demo", "Criar menu simples, objetivo de arena e condi\u00e7\u00e3o de vit\u00f3ria/derrota."),
    ]
    for i, (n, t, b) in enumerate(roadmap):
        x = 0.85 + i * 3.05
        add_card(slide, x, 1.75, 2.65, 3.05, f"{n}. {t}", b, [GOLD, CYAN, GREEN, BLUE][i], RGBColor(9, 18, 45))
    add_card(
        slide,
        1.2,
        5.45,
        10.9,
        0.95,
        "Mensagem final",
        "Mage j\u00e1 possui uma base jog\u00e1vel coerente: personagem, c\u00e2mera, mira, vida, UI, anima\u00e7\u00e3o e magia. O pr\u00f3ximo ganho est\u00e1 em fechar o ciclo de desafio.",
        CYAN,
        RGBColor(9, 18, 45),
    )
    add_footer(slide, "Mage | Pr\u00f3ximos passos")


def main():
    OUT_DIR.mkdir(exist_ok=True)
    if PPTX_PATH.exists() and not BACKUP_PATH.exists():
        shutil.copy2(PPTX_PATH, BACKUP_PATH)

    prepare_images()
    build_slides()

    prs.core_properties.title = "Mage - Apresentacao do Projeto"
    prs.core_properties.subject = "Apresentacao baseada no projeto Unity Mage"
    prs.core_properties.author = "Codex"
    prs.core_properties.comments = (
        "Gerado a partir de leitura local do projeto, Unity MCP e capturas reais da cena."
    )
    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
